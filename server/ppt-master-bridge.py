#!/usr/bin/env python3
"""
PPT Master Bridge - Python 脚本
从 stdin 读取 JSON 格式的 PPT 大纲，使用 python-pptx 生成精美可编辑的 PPTX 文件
设计理念参考 PPT Master (https://github.com/hugohe3/ppt-master)

用法: echo '{"title":"...","slides":[...]}' | python ppt-master-bridge.py output.pptx
"""

import sys
import json
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print(json.dumps({"error": "python-pptx not installed. Run: pip install python-pptx"}), file=sys.stderr)
    sys.exit(1)

# 配色方案 - 深海蓝科技风
COLORS = {
    'bg': RGBColor(0x0A, 0x0F, 0x1E),
    'bg_light': RGBColor(0x11, 0x1A, 0x2E),
    'accent': RGBColor(0x38, 0xBD, 0xF8),
    'accent_dark': RGBColor(0x02, 0x84, 0xC7),
    'text_primary': RGBColor(0xF1, 0xF5, 0xF9),
    'text_secondary': RGBColor(0x94, 0xA3, 0xB8),
    'text_muted': RGBColor(0x64, 0x74, 0x8B),
    'border': RGBColor(0x1E, 0x29, 0x3B),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
}

FONT_HEADING = 'Microsoft YaHei'
FONT_BODY = 'Microsoft YaHei'


def add_background(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color=None, line_color=None, line_width=None):
    """添加矩形装饰"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, x, y, w, h, color=None, line_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_ellipse(slide, x, y, w, h, color=None, line_color=None, line_width=None):
    """添加椭圆"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, font_size=14, color=None, bold=False, font_name=None, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    if color is None:
        color = COLORS['text_primary']
    if font_name is None:
        font_name = FONT_BODY
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def create_hero_slide(prs, slide_data, total_slides):
    """创建封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, COLORS['bg'])

    # 顶部 accent 线
    add_rect(slide, 0.5, 0.5, 2.5, 0.04, color=COLORS['accent'])

    # 主标题
    title = slide_data.get('title', 'Untitled')
    add_text(slide, 0.5, 2.0, 12.3, 1.0, title, font_size=40, color=COLORS['text_primary'], bold=True, font_name=FONT_HEADING)

    # 副标题
    subtitle = slide_data.get('subtitle', '')
    if subtitle:
        add_text(slide, 0.5, 3.1, 10, 0.5, subtitle, font_size=18, color=COLORS['text_secondary'])

    # 底部信息
    add_text(slide, 0.5, 5.5, 6, 0.3, f'共 {total_slides} 页 · AI 智能生成', font_size=11, color=COLORS['text_muted'])

    # 右下角装饰圆环
    add_ellipse(slide, 10.5, 4.8, 1.8, 1.8, line_color=COLORS['accent'], line_width=1.5)

    # 页码
    add_text(slide, 11.8, 0.35, 0.8, 0.3, '1', font_size=10, color=COLORS['text_muted'], alignment=PP_ALIGN.RIGHT)

    return slide


def create_text_slide(prs, slide_data, index, total_slides):
    """创建文字页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg'])

    # 左侧 accent 竖条
    add_rect(slide, 0.5, 0.6, 0.06, 0.5, color=COLORS['accent'])

    # 标题
    title = slide_data.get('title', f'Page {index + 1}')
    add_text(slide, 0.7, 0.55, 11.5, 0.55, title, font_size=26, color=COLORS['text_primary'], bold=True, font_name=FONT_HEADING)

    # 内容要点
    items = slide_data.get('content', [])[:6]
    for idx, item in enumerate(items):
        y = 1.4 + idx * 0.75
        # 序号圆圈
        add_ellipse(slide, 0.6, y + 0.04, 0.28, 0.28, color=COLORS['bg_light'], line_color=COLORS['border'], line_width=1)
        add_text(slide, 0.6, y + 0.04, 0.28, 0.28, str(idx + 1), font_size=10, color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
        # 文字
        add_text(slide, 1.1, y, 11, 0.45, str(item), font_size=15, color=COLORS['text_secondary'])

    if not items:
        add_text(slide, 1.1, 1.4, 11, 0.45, '暂无内容', font_size=14, color=COLORS['text_muted'])

    # 页码
    add_text(slide, 11.8, 0.35, 0.8, 0.3, str(index + 1), font_size=10, color=COLORS['text_muted'], alignment=PP_ALIGN.RIGHT)

    return slide


def create_cards_slide(prs, slide_data, index, total_slides):
    """创建卡片页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg'])

    # 标题
    title = slide_data.get('title', f'Cards {index + 1}')
    add_text(slide, 0.5, 0.55, 11.5, 0.55, title, font_size=26, color=COLORS['text_primary'], bold=True, font_name=FONT_HEADING)
    add_rect(slide, 0.5, 1.15, 1.0, 0.03, color=COLORS['accent'])

    # 卡片网格
    items = slide_data.get('items', [])[:4]
    for idx, item in enumerate(items):
        col = idx % 2
        row = idx // 2
        card_x = 0.5 + col * 6.15
        card_y = 1.5 + row * 2.4

        # 卡片背景
        add_rounded_rect(slide, card_x, card_y, 5.9, 2.1, color=COLORS['bg_light'], line_color=COLORS['border'])
        # 卡片顶部 accent 条
        add_rect(slide, card_x, card_y, 5.9, 0.04, color=COLORS['accent'])
        # 卡片标题
        add_text(slide, card_x + 0.25, card_y + 0.2, 5.4, 0.4, item.get('title', ''), font_size=16, color=COLORS['text_primary'], bold=True, font_name=FONT_HEADING)
        # 卡片描述
        add_text(slide, card_x + 0.25, card_y + 0.65, 5.4, 1.2, item.get('desc', ''), font_size=12, color=COLORS['text_secondary'])

    # 页码
    add_text(slide, 11.8, 0.35, 0.8, 0.3, str(index + 1), font_size=10, color=COLORS['text_muted'], alignment=PP_ALIGN.RIGHT)

    return slide


def create_closing_slide(prs, doc_title, index, total_slides):
    """创建结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['bg'])

    # 居中线
    add_rect(slide, 4.5, 2.2, 4.3, 0.04, color=COLORS['accent'])

    # 感谢文字
    add_text(slide, 0.5, 2.6, 12.3, 0.8, '感谢聆听', font_size=36, color=COLORS['text_primary'], bold=True, font_name=FONT_HEADING, alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 3.4, 12.3, 0.4, 'THANKS FOR WATCHING', font_size=12, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 4.2, 12.3, 0.3, 'Shrimp Workspace AI · 智能演示', font_size=11, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)

    # 装饰圆点
    for i in range(3):
        color = COLORS['accent'] if i == 1 else COLORS['border']
        add_ellipse(slide, 6.15 + (i - 1) * 0.5, 5.2, 0.12, 0.12, color=color)

    # 页码
    add_text(slide, 11.8, 0.35, 0.8, 0.3, str(index + 1), font_size=10, color=COLORS['text_muted'], alignment=PP_ALIGN.RIGHT)

    return slide


def generate_pptx(doc, output_path):
    """生成 PPTX 文件"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = doc.get('slides', [])
    total = len(slides)

    if total == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, COLORS['bg'])
        add_text(slide, 0.5, 3.0, 12.3, 0.8, '暂无内容', font_size=28, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)
    else:
        for index, slide_data in enumerate(slides):
            slide_type = slide_data.get('type', 'text')
            is_first = index == 0
            is_last = index == total - 1

            if slide_type == 'hero' or is_first:
                create_hero_slide(prs, slide_data, total)
            elif slide_type == 'cards':
                create_cards_slide(prs, slide_data, index, total)
            elif is_last and total > 2:
                create_closing_slide(prs, doc.get('title', ''), index, total)
            else:
                create_text_slide(prs, slide_data, index, total)

    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: python ppt-master-bridge.py <output_path> [--input <json_file>]"}), file=sys.stderr)
        sys.exit(1)

    output_path = sys.argv[1]

    # 支持 --input 参数从文件读取 JSON，否则从 stdin 读取
    input_file = None
    if '--input' in sys.argv:
        idx = sys.argv.index('--input')
        if idx + 1 < len(sys.argv):
            input_file = sys.argv[idx + 1]

    try:
        if input_file:
            with open(input_file, 'r', encoding='utf-8') as f:
                input_data = f.read()
        else:
            # Windows 下 stdin 默认编码可能不是 UTF-8
            try:
                sys.stdin.reconfigure(encoding='utf-8')
            except Exception:
                pass
            input_data = sys.stdin.read()
        doc = json.loads(input_data)
    except json.JSONDecodeError as e:
        print(json.dumps({"error": f"Invalid JSON input: {e}"}), file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(json.dumps({"error": f"Input error: {e}"}), file=sys.stderr)
        sys.exit(1)

    try:
        result_path = generate_pptx(doc, output_path)
        print(json.dumps({"success": True, "path": result_path}))
    except Exception as e:
        print(json.dumps({"error": str(e)}), file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()
